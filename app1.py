import os
import tempfile
import streamlit as st
import boto3
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_aws import BedrockEmbeddings, ChatBedrock
from langchain.prompts import PromptTemplate
from langchain.chains import RetrievalQA
import pickle

# ------------------------
# CONFIG
# ------------------------
BUCKET_NAME = "chat-bucket-rag-us"
INDEX_FAISS = "index/faiss_index"
INDEX_PKL = "index/faiss_store.pkl"

# AWS Clients
s3 = boto3.client("s3", region_name="us-east-1")  # specify your region
bedrock_client = boto3.client("bedrock-runtime", region_name="us-east-1")

# Embedding model
embeddings = BedrockEmbeddings(
    model_id="amazon.titan-embed-text-v2:0", client=bedrock_client
)

# ------------------------
# HELPERS
# ------------------------
def list_s3_files(bucket):
    """List all files in the S3 bucket root."""
    res = s3.list_objects_v2(Bucket=bucket)
    if "Contents" not in res:
        return []
    return [obj["Key"] for obj in res["Contents"] if obj["Key"].lower().endswith((".pdf",".docx",".pptx",".txt"))]

def download_s3_file(bucket, key):
    """Download a file from S3 to a temporary location."""
    tmp_file = tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(key)[1])
    s3.download_file(bucket, key, tmp_file.name)
    return tmp_file.name

def extract_text(file_path):
    ext = file_path.split(".")[-1].lower()
    text = ""
    if ext == "pdf":
        reader = PdfReader(file_path)
        for page in reader.pages:
            text += page.extract_text() or ""
    elif ext == "docx":
        doc = Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    elif ext == "pptx":
        pres = Presentation(file_path)
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    elif ext == "txt":
        with open(file_path, "r", encoding="utf-8") as f:
            text += f.read()
    return text

def build_index_from_s3():
    """Build FAISS index from all files in the S3 bucket."""
    files = list_s3_files(BUCKET_NAME)
    if not files:
        st.warning("No files found in S3 to build knowledge base.")
        return None

    all_texts = []
    for key in files:
        tmp_path = download_s3_file(BUCKET_NAME, key)
        all_texts.append(extract_text(tmp_path))

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    docs = text_splitter.create_documents(all_texts)
    db = FAISS.from_documents(docs, embeddings)

    # Save locally
    db.save_local("faiss_index")
    with open("faiss_index/index.pkl", "wb") as f:
        pickle.dump(db, f)

    # Upload to S3
    s3.upload_file("faiss_index/index.faiss", BUCKET_NAME, INDEX_FAISS)
    s3.upload_file("faiss_index/index.pkl", BUCKET_NAME, INDEX_PKL)

    return db

def load_index():
    """Load FAISS index from S3 if exists, else build from bucket files."""
    try:
        os.makedirs("faiss_index", exist_ok=True)
        s3.download_file(BUCKET_NAME, INDEX_FAISS, "faiss_index/index.faiss")
        s3.download_file(BUCKET_NAME, INDEX_PKL, "faiss_index/index.pkl")
        with open("faiss_index/index.pkl", "rb") as f:
            db = pickle.load(f)
        return FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
    except:
        return build_index_from_s3()

def process_new_files(uploaded_files, db):
    """Add new uploaded files to existing index and save to S3."""
    if not uploaded_files:
        return db

    all_texts = []
    for uploaded_file in uploaded_files:
        ext = uploaded_file.name.split(".")[-1].lower()
        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{ext}") as tmp:
            tmp.write(uploaded_file.read())
            tmp_path = tmp.name
        all_texts.append(extract_text(tmp_path))

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    docs = text_splitter.create_documents(all_texts)

    if db:
        db.add_documents(docs)
    else:
        db = FAISS.from_documents(docs, embeddings)

    # Save & upload
    db.save_local("faiss_index")
    with open("faiss_index/index.pkl", "wb") as f:
        pickle.dump(db, f)
    s3.upload_file("faiss_index/index.faiss", BUCKET_NAME, INDEX_FAISS)
    s3.upload_file("faiss_index/index.pkl", BUCKET_NAME, INDEX_PKL)

    return db

# ------------------------
# STREAMLIT APP
# ------------------------
st.title("📚 RAG Chatbot with AWS Bedrock + S3")

# Load or build index on startup
if "db" not in st.session_state:
    st.session_state.db = load_index()

# File uploader
uploaded_files = st.file_uploader(
    "Upload files (PDF, DOCX, PPTX, TXT)",
    type=["pdf", "docx", "pptx", "txt"],
    accept_multiple_files=True
)
if st.button("Build / Update Knowledge Base"):
    st.session_state.db = process_new_files(uploaded_files, st.session_state.db)
    st.success("Knowledge base updated!")

# Ask questions
if st.session_state.db:
    retriever = st.session_state.db.as_retriever()
    llm = ChatBedrock(model_id="meta.llama3-2-3b-instruct-v1:0", client=bedrock_client)
    qa_chain = RetrievalQA.from_chain_type(llm, retriever=retriever)

    query = st.text_input("Ask a question about your documents:")
    if query:
        answer = qa_chain.run(query)
        st.write("### Answer:")
        st.write(answer)
else:
    st.info("No knowledge base found. Please upload files or add files to S3.")
