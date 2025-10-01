import os
import tempfile
import streamlit as st
import boto3
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_aws import BedrockEmbeddings, ChatBedrock
from langchain.prompts import PromptTemplate
from langchain.chains import ConversationalRetrievalChain
from langchain.memory import ConversationBufferMemory

# ------------------------
# CONFIG
# ------------------------
BUCKET_NAME = "chat-bucket-rag-us"
INDEX_FAISS = "index/faiss_index"
AWS_REGION = "us-east-1"  # change if needed
os.environ["AWS_REGION"] = AWS_REGION

# AWS Clients
s3 = boto3.client("s3", region_name=AWS_REGION)
bedrock_client = boto3.client("bedrock-runtime", region_name=AWS_REGION)

# ------------------------
# HELPERS
# ------------------------
def load_embeddings():
    return BedrockEmbeddings(model_id="amazon.titan-embed-text-v2:0", client=bedrock_client)

def load_index():
    try:
        os.makedirs("faiss_index", exist_ok=True)
        s3.download_file(BUCKET_NAME, INDEX_FAISS, "faiss_index/index.faiss")
        embeddings = load_embeddings()
        db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
        return db
    except Exception as e:
        st.warning(f"No existing index found in S3: {e}")
        return None

def save_index(db):
    db.save_local("faiss_index")
    s3.upload_file("faiss_index/index.faiss", BUCKET_NAME, INDEX_FAISS)
    st.success("FAISS index saved to S3")

def extract_text_from_file(file_path, file_type):
    text = ""
    if file_type == "pdf":
        reader = PdfReader(file_path)
        for page in reader.pages:
            text += page.extract_text() + "\n"
    elif file_type == "docx":
        doc = Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    elif file_type == "pptx":
        pres = Presentation(file_path)
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    elif file_type == "txt":
        with open(file_path, "r", encoding="utf-8") as f:
            text += f.read()
    return text

def process_new_files(uploaded_files, db=None):
    all_texts = []
    for uploaded_file in uploaded_files:
        file_type = uploaded_file.name.split(".")[-1].lower()
        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_type}") as tmp_file:
            tmp_file.write(uploaded_file.read())
            tmp_path = tmp_file.name
        text = extract_text_from_file(tmp_path, file_type)
        if not text.strip():
            st.warning(f"No text extracted from {uploaded_file.name}")
        all_texts.append(text)

    if not all_texts:
        return db

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    docs = text_splitter.create_documents(all_texts)

    if db:
        db.add_documents(docs)
    else:
        embeddings = load_embeddings()
        db = FAISS.from_documents(docs, embeddings)

    save_index(db)
    return db

# ------------------------
# STREAMLIT APP
# ------------------------
st.title("RAG Chatbot with AWS Bedrock + Context Memory")

# Load knowledge base
if "db" not in st.session_state:
    st.session_state.db = load_index()

# Init chat history
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []

# Init LLM + chain
if "qa_chain" not in st.session_state and st.session_state.db:
    retriever = st.session_state.db.as_retriever()
    llm = ChatBedrock(
        model_id="arn:aws:bedrock:us-east-1:703139106624:inference-profile/us.meta.llama3-2-3b-instruct-v1:0",
        model_kwargs={},
        provider="meta"
    )

    memory = ConversationBufferMemory(memory_key="chat_history", return_messages=True)

    # Custom prompt that includes chat history
    prompt_template = """You are a helpful assistant. 
Here is the conversation so far:
{chat_history}

User question: {question}
Answer:"""

    qa_prompt = PromptTemplate(
        input_variables=["chat_history", "question"],
        template=prompt_template
    )

    st.session_state.qa_chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=retriever,
        memory=memory,
        combine_docs_chain_kwargs={"prompt": qa_prompt}
    )

# Upload files
uploaded_files = st.file_uploader(
    "Upload files (PDF, DOCX, PPTX, TXT)",
    type=["pdf", "docx", "pptx", "txt"],
    accept_multiple_files=True
)

if st.button("Build / Update Knowledge Base"):
    st.session_state.db = process_new_files(uploaded_files, st.session_state.db)
    if st.session_state.db:
        retriever = st.session_state.db.as_retriever()
        llm = ChatBedrock(
            model_id="arn:aws:bedrock:us-east-1:703139106624:inference-profile/us.meta.llama3-2-3b-instruct-v1:0",
            model_kwargs={},
            provider="meta"
        )
        memory = ConversationBufferMemory(memory_key="chat_history", return_messages=True)

        prompt_template = """You are a helpful assistant. 
Here is the conversation so far:
{chat_history}

User question: {question}
Answer:"""

        qa_prompt = PromptTemplate(
            input_variables=["chat_history", "question"],
            template=prompt_template
        )

        st.session_state.qa_chain = ConversationalRetrievalChain.from_llm(
            llm=llm,
            retriever=retriever,
            memory=memory,
            combine_docs_chain_kwargs={"prompt": qa_prompt}
        )

# Chat interface
if st.session_state.db and "qa_chain" in st.session_state:
    query = st.text_input("Ask a question about your documents:")
    if query:
        with st.spinner("Generating answer..."):
            result = st.session_state.qa_chain({"question": query})
            answer = result["answer"]
            st.session_state.chat_history.append((query, answer))

    # Display chat history
    if st.session_state.chat_history:
        st.write("### Conversation History")
        for q, a in st.session_state.chat_history:
            st.write(f"**You:** {q}")
            st.write(f"**Bot:** {a}")
else:
    st.info("No knowledge base found. Please upload files and click 'Build / Update Knowledge Base'.")
