import os
import tempfile
import streamlit as st
import boto3
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_aws import BedrockEmbeddings
from langchain.prompts import PromptTemplate
from langchain.chains import RetrievalQA
from langchain_aws import ChatBedrock
import pickle

# ------------------------
# CONFIG
# ------------------------
BUCKET_NAME = "chat-bucket-rag-us"
INDEX_FAISS = "index/faiss_index"
INDEX_PKL = "index/faiss_store.pkl"

# AWS Clients (no keys needed if EC2 has IAM role)
s3 = boto3.client("s3")
bedrock_client = boto3.client("bedrock-runtime",region_name="us-east-1")

# Embedding Model
embeddings = BedrockEmbeddings(
    model_id="amazon.titan-embed-text-v1", client=bedrock_client
)

# ------------------------
# HELPERS
# ------------------------
def load_from_s3():
    try:
        os.makedirs("faiss_index", exist_ok=True)
        s3.download_file(BUCKET_NAME, INDEX_FAISS, "faiss_index/index.faiss")
        s3.download_file(BUCKET_NAME, INDEX_PKL, "faiss_index/index.pkl")
        with open("faiss_index/index.pkl", "rb") as f:
            store = pickle.load(f)
        return FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
    except Exception as e:
        st.warning(f"No existing index found in S3: {e}")
        return None

def save_to_s3(db):
    db.save_local("faiss_index")
    with open("faiss_index/index.pkl", "wb") as f:
        pickle.dump(db, f)
    s3.upload_file("faiss_index/index.faiss", BUCKET_NAME, INDEX_FAISS)
    s3.upload_file("faiss_index/index.pkl", BUCKET_NAME, INDEX_PKL)

def extract_text_from_file(file_path, file_type):
    text = ""
    if file_type == "pdf":
        reader = PdfReader(file_path)
        for page in reader.pages:
            text += page.extract_text() + "\n"
    elif file_type == "docx":
        doc = Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    elif file_type == "pptx":
        pres = Presentation(file_path)
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    elif file_type == "txt":
        with open(file_path, "r", encoding="utf-8") as f:
            text += f.read()
    return text

def process_new_files(uploaded_files, db=None):
    all_texts = []
    for uploaded_file in uploaded_files:
        file_type = uploaded_file.name.split(".")[-1].lower()
        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_type}") as tmp_file:
            tmp_file.write(uploaded_file.read())
            tmp_path = tmp_file.name
        text = extract_text_from_file(tmp_path, file_type)
        all_texts.append(text)

    if not all_texts:
        return db

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    docs = text_splitter.create_documents(all_texts)

    if db:
        db.add_documents(docs)
    else:
        db = FAISS.from_documents(docs, embeddings)

    save_to_s3(db)
    return db

# ------------------------
# STREAMLIT APP
# ------------------------
def main():
    st.title("📚 RAG Chatbot with AWS Bedrock + S3")

    # Step 1: Load knowledge base
    if "db" not in st.session_state:
        st.session_state.db = load_from_s3()

    uploaded_files = st.file_uploader("Upload files (PDF, DOCX, PPTX, TXT)", 
                                      type=["pdf", "docx", "pptx", "txt"], 
                                      accept_multiple_files=True)
    if st.button("Build / Update Knowledge Base"):
        st.session_state.db = process_new_files(uploaded_files, st.session_state.db)
        st.success("Knowledge base updated!")

    # Step 2: Ask questions
    if st.session_state.db:
        retriever = st.session_state.db.as_retriever()
        llm = ChatBedrock(model_id="anthropic.claude-v2:1", client=bedrock_client)
        qa_chain = RetrievalQA.from_chain_type(llm, retriever=retriever)

        query = st.text_input("Ask a question about your documents:")
        if query:
            answer = qa_chain.run(query)
            st.write("### Answer:")
            st.write(answer)
    else:
        st.info("No knowledge base found. Please upload files and click 'Build Knowledge Base'.")

if __name__ == "__main__":
    main()
